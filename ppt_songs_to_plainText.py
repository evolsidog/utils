import time
from glob import glob
from pptx import Presentation
import os
import win32com.client

"""
Cambiar el texto de "INPUT_DIR" y de "OUTPUT_DIR"
"""

START_TIME = time.time()

# Directorio con archivos de PowerPoint
INPUT_DIR = "./power_point_canciones"
INPUT_DIR = os.path.join(os.path.abspath(os.path.dirname(__file__)), INPUT_DIR)
# Directorio donde guardar los archivos de texto
OUTPUT_DIR = "./texto_plano_canciones"
OUTPUT_DIR = os.path.join(os.path.abspath(os.path.dirname(__file__)), OUTPUT_DIR)
os.makedirs(OUTPUT_DIR, exist_ok=True)

ARCHIVOS_ERROR_CONVERSION_PPT = []
ARCHIVOS_ERROR_EXTRACCION_TEXTO = []

def convert_ppt_to_pptx(ppt_path):
    """
    Función para convertir .ppt a .pptx
    """
    print(f"Convirtiendo ppt en pptx: {ppt_path}")

    # Aplicación de PowerPoint (necesario para convertir .ppt a .pptx)
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
    pptx_path = ppt_path.replace(".ppt", "_antiguo.pptx") # hay powerpoints con el mismo nombre y distinto formato
    presentation.SaveAs(pptx_path, FileFormat=24)  # 24 es el formato para .pptx
    presentation.Close()
    powerpoint.Quit()
    
    print(f"ppt convertida en pptx: {pptx_path}")
    return pptx_path


def process_pptx(pptx_path, output_file):
    """
    Función para procesar archivos PPTX y extraer texto
    """
    # Cargamos power point
    prs = Presentation(pptx_path)
    
    # Abrimos fichero de texto plano donde vamos a ir guardando el resultado
    with open(output_file, "w", encoding="utf-8") as f:
        # Iterar sobre cada diapositiva
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    # Limpiar texto de saltos de línea innecesarios
                    cleaned_text = paragraph.text.strip()
                    if cleaned_text:  # Solo añadir si no está vacío
                        slide_text.append(cleaned_text)
             # Escribir el texto de la diapositiva y añadir un salto de línea
            if slide_text:
                f.write("\n".join(slide_text) + "\n\n")
    
    print(f"Texto extraído en {output_file}")


# 1. Buscar todos los archivos de power point en el directorio de entrada y subdirectorios
pptx_files = glob(os.path.join(INPUT_DIR, '**', '*.pptx'), recursive=True)
# filenames = [filename for filename in os.listdir(INPUT_DIR) if filename.endswith(".pptx")]
print(f"Numero de pptx formato nuevo: {len(pptx_files)}")
ppt_files = glob(os.path.join(INPUT_DIR, '**', '*.ppt'), recursive=True)
print(f"Numero de ppt antiguos: {len(ppt_files)}")

# 2. Convertimos los ppt antiguos al nuevo formato pptx
ppts_converted = []
for ppt_path in ppt_files:
    try:
        # Convertir .ppt a .pptx
        ppt_converted = convert_ppt_to_pptx(ppt_path=ppt_path)
        ppts_converted.append(ppt_converted) 
    except Exception as e:
        print(f"Error convirtiendo formato antiguo: {ppt_path}")
        parent_directory = os.path.basename(os.path.dirname(ppt_path))
        filename = os.path.basename(ppt_path)
        ARCHIVOS_ERROR_CONVERSION_PPT.append(os.path.join(parent_directory, filename))    

print(f"Numero de ppt convertidos: {len(ppts_converted)}")

pptx_files = pptx_files + ppts_converted
print(f"Numero de pptx formato nuevo tras conversión: {len(pptx_files)}")

print(f"Tiempo total de conversión ppt antiguos: {time.time() - START_TIME:.4f} segundos")

# 3. Extraemos el texto de cada power point
for pptx in pptx_files:
    try:
        print(f"Leyendo power point {pptx}")

        pptx_path = os.path.join(INPUT_DIR, pptx)
        
        # Nombre base para el fichero de texto plano de salida
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        output_file = os.path.join(OUTPUT_DIR, f"{base_name}.txt")

        process_pptx(pptx_path=pptx_path, output_file=output_file)
    except Exception as e:
        print(f"Error extrayendo texto de powerpoint: {pptx}")
        parent_directory = os.path.basename(os.path.dirname(pptx_path))
        filename = os.path.basename(pptx_path)
        ARCHIVOS_ERROR_EXTRACCION_TEXTO.append(os.path.join(parent_directory, filename)) 


print("Extracción completada para todos los archivos.")
print(f"Numero de ppt con error al convertir a powerpoint moderno: {len(ARCHIVOS_ERROR_CONVERSION_PPT)}")
print(ARCHIVOS_ERROR_CONVERSION_PPT)
print(f"Numero de ppt con error al extraer texto: {len(ARCHIVOS_ERROR_EXTRACCION_TEXTO)}")
print(ARCHIVOS_ERROR_EXTRACCION_TEXTO)
print(f"Tiempo total del proceso: {time.time() - START_TIME:.4f} segundos")
